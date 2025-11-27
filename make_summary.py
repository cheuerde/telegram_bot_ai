"""
Document Summarization Module

Supports: PDF, DOCX, PPTX, TXT, URL
Uses OpenAI GPT-4o-mini for summarization
"""

import re
import io
import os
import logging
from typing import List, Tuple

import docx
import pptx
import requests
from bs4 import BeautifulSoup
from openai import OpenAI
from pdfminer.converter import TextConverter
from pdfminer.pdfinterp import PDFPageInterpreter, PDFResourceManager
from pdfminer.pdfpage import PDFPage
from pdfminer.layout import LAParams

logger = logging.getLogger(__name__)

# PDF extraction parameters
laparams = LAParams()
laparams.char_margin = 1
laparams.word_margin = 2

# Lazy-initialized OpenAI client
_openai_client = None


def get_openai_client():
    """Get or create OpenAI client (lazy initialization)."""
    global _openai_client
    if _openai_client is None:
        _openai_client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
    return _openai_client


def extract_text_by_page(pdf_path: str):
    """Extract text from PDF page by page."""
    with open(pdf_path, 'rb') as fh:
        for page in PDFPage.get_pages(fh, caching=True, check_extractable=True):
            resource_manager = PDFResourceManager()
            fake_file_handle = io.StringIO()
            converter = TextConverter(resource_manager, fake_file_handle, laparams=laparams)
            page_interpreter = PDFPageInterpreter(resource_manager, converter)
            page_interpreter.process_page(page)

            text = fake_file_handle.getvalue()
            yield text

            converter.close()
            fake_file_handle.close()


def extract_text(pdf_path: str) -> List[str]:
    """Extract and chunk text from PDF."""
    text = ""
    for page in extract_text_by_page(pdf_path):
        text += page
    return extract_chapters(text)


def extract_chapters(text: str, max_words: int = 1000) -> List[str]:
    """Split text into chunks of max_words."""
    words = text.split()
    chunks = []
    chunk = ""

    for word in words:
        if len(chunk.split()) + 1 <= max_words:
            chunk += " " + word
        else:
            chunks.append(chunk.strip())
            chunk = word

    if chunk.strip():
        chunks.append(chunk.strip())

    return chunks


def create_summary(text: str, max_tokens: int = 100, prompt_prefix: str = '') -> str:
    """Create summary using GPT-4o-mini."""
    try:
        prompt = prompt_prefix + text

        response = get_openai_client().chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert summarizer. Create concise, informative summaries. When asked for bullet points, use \\item format for LaTeX compatibility."
                },
                {"role": "user", "content": prompt}
            ],
            max_tokens=max_tokens,
            temperature=0.7
        )

        summary = response.choices[0].message.content.strip()
        summary = summary.rstrip(',').rstrip('.')
        return summary

    except Exception as e:
        logger.error(f"Summary creation error: {e}")
        return f"Error creating summary: {str(e)}"


def generate_summaries(chapters: List[str], min_words_summary: int = 20) -> List[str]:
    """Generate summaries for each chapter."""
    summaries = []
    prompt = (
        'Create a brief summary from the input text in bullet points. '
        'Start every bullet point with "\\item ". Do not output incomplete sentences. '
        'This is the input text: '
    )

    for chapter in chapters:
        if not chapter.strip():
            continue
        summary = create_summary(chapter, max_tokens=100, prompt_prefix=prompt)
        summaries.append(summary)

    # Remove very short final summaries
    if summaries and len(summaries[-1].split()) <= min_words_summary:
        summaries = summaries[:-1]

    if summaries:
        summaries[-1] += '.'

    return summaries


def summarize_pdf(summaries: List[str], overall_summary: str, file_out: str) -> Tuple[str, str]:
    """Create LaTeX PDF with summaries."""
    out_file_raw, _ = os.path.splitext(file_out)
    pdf_path = os.path.dirname(file_out)
    tex_file = out_file_raw + '.tex'

    with open(tex_file, 'w', encoding='utf-8') as f:
        f.write(r'\documentclass[12pt]{article}' + '\n')
        f.write(r'\usepackage[utf8]{inputenc}' + '\n')
        f.write(r'\usepackage{amsmath}' + '\n')
        f.write(r'\usepackage{amsfonts}' + '\n')
        f.write(r'\usepackage{amssymb}' + '\n')
        f.write(r'\usepackage{graphicx}' + '\n')
        f.write(r'\begin{document}' + '\n')

        f.write(r'\section*{Number of Chapters: ' + str(len(summaries)) + '}' + '\n')
        f.write(r'\subsection*{Overall Summary}' + '\n')
        f.write(overall_summary + '\n')
        f.write(r'\newpage' + '\n')
        f.write(r'\section*{Chapter Summaries}' + '\n')

        for i, summary in enumerate(summaries):
            f.write(r'\subsection*{Chapter ' + str(i + 1) + ' Summary}' + '\n')
            f.write(r'\begin{itemize}' + '\n')
            f.write(summary.replace("\\item", "\\item ") + '\n')
            f.write(r'\end{itemize}' + '\n')

        f.write(r'\end{document}' + '\n')

    # Compile LaTeX to PDF
    command = f'pdflatex -interaction=nonstopmode -output-directory={pdf_path} {tex_file} > /dev/null 2>&1'
    os.system(command)

    return file_out, overall_summary


def pdf_to_summary(file_in: str, file_out: str) -> Tuple[str, str]:
    """Convert PDF to summary PDF."""
    try:
        chapters = extract_text(file_in)
        if not chapters:
            return ("Error", "Could not extract text from PDF")

        summaries = generate_summaries(chapters, min_words_summary=10)
        if not summaries:
            return ("Error", "Could not generate summaries")

        combined_text = " ".join(summaries).replace('\\item', '')
        overall_summary = create_summary(
            text=combined_text,
            max_tokens=400,
            prompt_prefix='From the given text, generate a concise overall summary: '
        )

        return summarize_pdf(summaries, overall_summary, file_out)

    except Exception as e:
        logger.error(f"PDF summary error: {e}")
        return ("Error", f"Failed to process PDF: {str(e)}")


def txt_to_summary(file_in: str, file_out: str) -> Tuple[str, str]:
    """Convert text file to summary PDF."""
    try:
        with open(file_in, 'r', encoding='utf-8') as file:
            text = file.read()

        chapters = extract_chapters(text, max_words=1000)
        if not chapters:
            return ("Error", "Could not extract text from file")

        summaries = generate_summaries(chapters, min_words_summary=10)
        if not summaries:
            return ("Error", "Could not generate summaries")

        combined_text = " ".join(summaries).replace('\\item', '')
        overall_summary = create_summary(
            text=combined_text,
            max_tokens=400,
            prompt_prefix='From the given text, generate a concise overall summary: '
        )

        return summarize_pdf(summaries, overall_summary, file_out)

    except Exception as e:
        logger.error(f"TXT summary error: {e}")
        return ("Error", f"Failed to process text file: {str(e)}")


def docx_to_summary(file_in: str, file_out: str) -> Tuple[str, str]:
    """Convert Word document to summary PDF."""
    try:
        document = docx.Document(file_in)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"

        chapters = extract_chapters(text, max_words=1000)
        if not chapters:
            return ("Error", "Could not extract text from document")

        summaries = generate_summaries(chapters, min_words_summary=10)
        if not summaries:
            return ("Error", "Could not generate summaries")

        combined_text = " ".join(summaries).replace('\\item', '')
        overall_summary = create_summary(
            text=combined_text,
            max_tokens=400,
            prompt_prefix='From the given text, generate a concise overall summary: '
        )

        return summarize_pdf(summaries, overall_summary, file_out)

    except Exception as e:
        logger.error(f"DOCX summary error: {e}")
        return ("Error", f"Failed to process Word document: {str(e)}")


def pptx_to_summary(file_in: str, file_out: str) -> Tuple[str, str]:
    """Convert PowerPoint to summary PDF."""
    try:
        presentation = pptx.Presentation(file_in)
        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"

        chapters = extract_chapters(text, max_words=1000)
        if not chapters:
            return ("Error", "Could not extract text from presentation")

        summaries = generate_summaries(chapters, min_words_summary=10)
        if not summaries:
            return ("Error", "Could not generate summaries")

        combined_text = " ".join(summaries).replace('\\item', '')
        overall_summary = create_summary(
            text=combined_text,
            max_tokens=400,
            prompt_prefix='From the given text, generate a concise overall summary: '
        )

        return summarize_pdf(summaries, overall_summary, file_out)

    except Exception as e:
        logger.error(f"PPTX summary error: {e}")
        return ("Error", f"Failed to process PowerPoint: {str(e)}")


def url_to_summary(url_in: str, file_out: str) -> Tuple[str, str]:
    """Convert URL content to summary PDF."""
    try:
        response = requests.get(url_in, timeout=30)
        response.raise_for_status()

        soup = BeautifulSoup(response.text, 'html.parser')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        text = soup.get_text()

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)

        chapters = extract_chapters(text, max_words=1000)
        if not chapters:
            return ("Error", "Could not extract text from URL")

        summaries = generate_summaries(chapters, min_words_summary=10)
        if not summaries:
            return ("Error", "Could not generate summaries")

        combined_text = " ".join(summaries).replace('\\item', '')
        overall_summary = create_summary(
            text=combined_text,
            max_tokens=400,
            prompt_prefix='From the given text, generate a concise overall summary: '
        )

        return summarize_pdf(summaries, overall_summary, file_out)

    except requests.RequestException as e:
        logger.error(f"URL fetch error: {e}")
        return ("Error", f"Failed to fetch URL: {str(e)}")
    except Exception as e:
        logger.error(f"URL summary error: {e}")
        return ("Error", f"Failed to process URL: {str(e)}")
